from flask import Flask, request, jsonify, Response, stream_with_context, send_file
from flask_cors import CORS
from flask_jwt_extended import JWTManager, create_access_token, jwt_required, get_jwt_identity
from pymongo import MongoClient
from bson import ObjectId
from datetime import datetime, timedelta
import requests
import json
import bcrypt
import os
from dotenv import load_dotenv
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
import re
import uuid
import PyPDF2
import base64
from PIL import Image, ImageEnhance, ImageFilter
import pytesseract
from io import BytesIO
import sys
import traceback
import platform
import subprocess
from pathlib import Path
from werkzeug.utils import secure_filename
import shutil
import time

# Optional imports
try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx not installed")

try:
    import fitz
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("⚠️ PyMuPDF not installed")

try:
    import pythoncom
    import win32com.client
    WIN32COM_AVAILABLE = True
except ImportError:
    WIN32COM_AVAILABLE = False
    print("⚠️ win32com not available")

# Load environment variables
load_dotenv()

app = Flask(__name__)

# Configure CORS - Allow all origins for development
CORS(app, origins=[
    "http://localhost:3000",
    "http://localhost:3003",
    "http://localhost:3005",
    "http://127.0.0.1:3000",
    "http://127.0.0.1:3003",
    "http://127.0.0.1:3005"
], supports_credentials=True, allow_headers=["Content-Type", "Authorization", "Accept"], methods=["GET", "POST", "PUT", "DELETE", "OPTIONS", "PATCH"])

# JWT Configuration
app.config["JWT_SECRET_KEY"] = os.getenv("JWT_SECRET_KEY", "dev-secret-key-change-this-in-production")
app.config["JWT_ACCESS_TOKEN_EXPIRES"] = timedelta(days=7)
jwt = JWTManager(app)

# MongoDB Connection
MONGODB_URI = os.getenv("MONGODB_URI", "mongodb://localhost:27017/")
try:
    client = MongoClient(MONGODB_URI, serverSelectionTimeoutMS=5000)
    client.admin.command('ping')
    print("✅ Connected to MongoDB successfully!")
    db = client["llama_gpt"]
    users_col = db["users"]
    chats_col = db["chats"]
    messages_col = db["messages"]
    files_col = db["files"]
    documents_col = db["documents"]
    converted_files_col = db["converted_files"]
    
    users_col.create_index("username", unique=True)
    chats_col.create_index([("user_id", 1), ("updated_at", -1)])
    messages_col.create_index([("chat_id", 1), ("created_at", 1)])
    files_col.create_index([("user_id", 1), ("uploaded_at", -1)])
    documents_col.create_index([("user_id", 1), ("created_at", -1)])
    converted_files_col.create_index([("user_id", 1), ("created_at", -1)])
    
except Exception as e:
    print(f"❌ MongoDB connection error: {e}")
    sys.exit(1)

# Ollama Configuration
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://localhost:11434")
DEFAULT_MODEL = "llama3.2"

# Configure upload folders
UPLOAD_FOLDER = 'uploads'
CONVERTED_FOLDER = 'converted_files'
DOCUMENTS_FOLDER = 'generated_documents'

for folder in [UPLOAD_FOLDER, CONVERTED_FOLDER, DOCUMENTS_FOLDER]:
    if not os.path.exists(folder):
        os.makedirs(folder)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['CONVERTED_FOLDER'] = CONVERTED_FOLDER
app.config['DOCUMENTS_FOLDER'] = DOCUMENTS_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024

system = platform.system()
home_dir = str(Path.home())

def serialize(doc):
    if doc and "_id" in doc:
        doc["_id"] = str(doc["_id"])
    return doc

def _build_cors_preflight_response():
    response = jsonify({"status": "ok"})
    response.headers.add("Access-Control-Allow-Origin", "*")
    response.headers.add("Access-Control-Allow-Headers", "Content-Type, Authorization, Accept")
    response.headers.add("Access-Control-Allow-Methods", "GET, POST, PUT, DELETE, OPTIONS, PATCH")
    response.headers.add("Access-Control-Allow-Credentials", "true")
    return response

def _corsify_actual_response(response):
    response.headers.add("Access-Control-Allow-Origin", "*")
    response.headers.add("Access-Control-Allow-Credentials", "true")
    return response

# ==================== FILE CONVERTER FUNCTIONS ====================

def convert_pdf_to_word(pdf_path, output_path):
    """Convert PDF to Word document"""
    try:
        doc = Document()
        
        section = doc.sections[0]
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
        
        style = doc.styles['Normal']
        style.font.name = 'Arial'
        style.font.size = Pt(11)
        
        with open(pdf_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text:
                    if page_num > 0:
                        doc.add_page_break()
                    
                    header_para = doc.add_paragraph()
                    header_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    header_run = header_para.add_run(f"Page {page_num + 1}")
                    header_run.font.size = Pt(10)
                    header_run.font.italic = True
                    
                    lines = text.split('\n')
                    for line in lines:
                        if line.strip():
                            p = doc.add_paragraph(line)
                            p.paragraph_format.line_spacing = 1.5
                            p.paragraph_format.space_after = Pt(6)
        
        doc.save(output_path)
        print(f"✅ PDF to Word conversion successful")
        return True
        
    except Exception as e:
        print(f"PDF to Word conversion error: {e}")
        traceback.print_exc()
        return False

def convert_word_to_pdf(word_path, output_path):
    """Convert Word document to PDF - FIXED VERSION"""
    try:
        print(f"🔄 Converting Word to PDF: {word_path} -> {output_path}")
        
        if not os.path.exists(word_path):
            print(f"❌ Input file does not exist: {word_path}")
            return False
        
        # Method 1: Try using LibreOffice (best cross-platform solution)
        libreoffice_paths = []
        if system == "Windows":
            libreoffice_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                r"C:\Program Files\LibreOffice\program\soffice.bin",
            ]
        elif system == "Darwin":  # macOS
            libreoffice_paths = [
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "/usr/local/bin/soffice"
            ]
        else:  # Linux
            libreoffice_paths = ["libreoffice", "soffice", "/usr/bin/libreoffice"]
        
        for lo_path in libreoffice_paths:
            try:
                # Check if LibreOffice exists
                if system == "Windows":
                    if not os.path.exists(lo_path):
                        continue
                else:
                    # For Linux/Mac, check if command exists
                    check = subprocess.run(["which", lo_path], capture_output=True, text=True)
                    if check.returncode != 0:
                        continue
                
                print(f"📝 Found LibreOffice at: {lo_path}")
                
                # Convert using LibreOffice
                cmd = [
                    lo_path,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", os.path.dirname(output_path),
                    word_path
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                
                if result.returncode == 0:
                    expected_pdf = os.path.join(
                        os.path.dirname(output_path),
                        os.path.splitext(os.path.basename(word_path))[0] + ".pdf"
                    )
                    
                    if os.path.exists(expected_pdf):
                        if expected_pdf != output_path:
                            shutil.move(expected_pdf, output_path)
                        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                            print(f"✅ Word to PDF conversion successful (LibreOffice)")
                            return True
            except Exception as e:
                print(f"LibreOffice at {lo_path} failed: {e}")
                continue
        
        # Method 2: On Windows, try COM automation
        if system == "Windows" and WIN32COM_AVAILABLE:
            try:
                print("Attempting Windows COM automation...")
                pythoncom.CoInitialize()
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                word.DisplayAlerts = False
                
                doc = word.Documents.Open(os.path.abspath(word_path))
                doc.SaveAs(os.path.abspath(output_path), FileFormat=17)
                doc.Close()
                word.Quit()
                pythoncom.CoUninitialize()
                
                if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    print(f"✅ Word to PDF conversion successful (Windows COM)")
                    return True
            except Exception as e:
                print(f"Windows COM conversion failed: {e}")
                traceback.print_exc()
        
        # Method 3: Create a simple PDF using extracted text
        try:
            print("Attempting text extraction fallback...")
            doc = Document(word_path)
            text_content = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # Create a simple text-based PDF
            temp_txt = output_path.replace('.pdf', '_temp.txt')
            with open(temp_txt, 'w', encoding='utf-8') as f:
                f.write('\n\n'.join(text_content))
            
            # Try to convert text to PDF
            success = convert_text_to_pdf(temp_txt, output_path)
            
            if os.path.exists(temp_txt):
                os.remove(temp_txt)
            
            if success:
                print(f"✅ Word to PDF conversion successful (text fallback)")
                return True
        except Exception as e:
            print(f"Text fallback failed: {e}")
        
        print(f"❌ All Word to PDF conversion methods failed")
        return False
        
    except Exception as e:
        print(f"Word to PDF conversion error: {e}")
        traceback.print_exc()
        return False

def convert_ppt_to_pdf(ppt_path, output_path):
    """Convert PowerPoint to PDF"""
    try:
        print(f"🔄 Converting PPT to PDF: {ppt_path} -> {output_path}")
        
        if not os.path.exists(ppt_path):
            print(f"❌ Input file does not exist: {ppt_path}")
            return False
        
        # Try LibreOffice
        libreoffice_paths = []
        if system == "Windows":
            libreoffice_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            ]
        elif system == "Darwin":
            libreoffice_paths = ["/Applications/LibreOffice.app/Contents/MacOS/soffice"]
        else:
            libreoffice_paths = ["libreoffice", "soffice"]
        
        for lo_path in libreoffice_paths:
            try:
                if system == "Windows":
                    if not os.path.exists(lo_path):
                        continue
                else:
                    check = subprocess.run(["which", lo_path], capture_output=True, text=True)
                    if check.returncode != 0:
                        continue
                
                cmd = [
                    lo_path,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", os.path.dirname(output_path),
                    ppt_path
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                
                if result.returncode == 0:
                    expected_pdf = os.path.join(
                        os.path.dirname(output_path),
                        os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf"
                    )
                    
                    if os.path.exists(expected_pdf):
                        if expected_pdf != output_path:
                            shutil.move(expected_pdf, output_path)
                        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                            print(f"✅ PPT to PDF conversion successful (LibreOffice)")
                            return True
            except Exception as e:
                continue
        
        # Try Windows COM
        if system == "Windows" and WIN32COM_AVAILABLE:
            try:
                pythoncom.CoInitialize()
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = False
                powerpoint.DisplayAlerts = False
                
                presentation = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
                presentation.SaveAs(os.path.abspath(output_path), 32)
                presentation.Close()
                powerpoint.Quit()
                pythoncom.CoUninitialize()
                
                if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    print(f"✅ PPT to PDF conversion successful (Windows COM)")
                    return True
            except Exception as e:
                print(f"Windows COM conversion failed: {e}")
        
        print(f"❌ All PPT to PDF conversion methods failed")
        return False
        
    except Exception as e:
        print(f"PPT to PDF conversion error: {e}")
        traceback.print_exc()
        return False

def convert_pdf_to_ppt(pdf_path, output_path):
    """Convert PDF to PowerPoint"""
    try:
        if not PPTX_AVAILABLE:
            return False
        
        prs = Presentation()
        
        with open(pdf_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                
                if text:
                    if page_num == 0:
                        slide = prs.slides.add_slide(prs.slide_layouts[0])
                        title = slide.shapes.title
                        subtitle = slide.placeholders[1]
                        if title:
                            title.text = f"PDF Document"
                        if subtitle:
                            first_line = text.split('\n')[0] if text else ""
                            subtitle.text = first_line[:100] if first_line else "Converted from PDF"
                    else:
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        title_box = slide.shapes.add_textbox(
                            PptxInches(0.5), PptxInches(0.5), PptxInches(9), PptxInches(1)
                        )
                        title_frame = title_box.text_frame
                        title_frame.text = f"Page {page_num + 1}"
                        title_frame.paragraphs[0].font.size = Pt(24)
                        title_frame.paragraphs[0].font.bold = True
                        
                        content_box = slide.shapes.add_textbox(
                            PptxInches(0.5), PptxInches(1.5), PptxInches(9), PptxInches(5.5)
                        )
                        content_frame = content_box.text_frame
                        content_frame.text = text[:2000]
        
        prs.save(output_path)
        print(f"✅ PDF to PPT conversion successful")
        return True
        
    except Exception as e:
        print(f"PDF to PPT conversion error: {e}")
        return False

def convert_image_to_pdf(image_path, output_path):
    """Convert image to PDF"""
    try:
        image = Image.open(image_path)
        if image.mode == 'RGBA':
            image = image.convert('RGB')
        image.save(output_path, 'PDF')
        return True
    except Exception as e:
        print(f"Image to PDF conversion error: {e}")
        return False

def convert_pdf_to_image(pdf_path, output_path, format='png'):
    """Convert PDF to image"""
    try:
        if not PYMUPDF_AVAILABLE:
            return False
        
        pdf_doc = fitz.open(pdf_path)
        page = pdf_doc[0]
        pix = page.get_pixmap()
        pix.save(output_path)
        pdf_doc.close()
        return True
    except Exception as e:
        print(f"PDF to Image conversion error: {e}")
        return False

def convert_text_to_pdf(text_path, output_path):
    """Convert text file to PDF"""
    try:
        with open(text_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Try using fpdf if available
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=11)
            
            for line in content.split('\n'):
                if line.strip():
                    pdf.multi_cell(0, 6, line[:500])
                    pdf.ln(4)
            
            pdf.output(output_path)
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                return True
        except ImportError:
            pass
        
        # Fallback: Create a Word document and convert
        doc = Document()
        for line in content.split('\n'):
            if line.strip():
                doc.add_paragraph(line)
        
        temp_word = output_path.replace('.pdf', '_temp.docx')
        doc.save(temp_word)
        success = convert_word_to_pdf(temp_word, output_path)
        
        if os.path.exists(temp_word):
            os.remove(temp_word)
        
        return success
            
    except Exception as e:
        print(f"Text to PDF conversion error: {e}")
        return False

def convert_text_to_word(text_path, output_path):
    """Convert text file to Word"""
    try:
        with open(text_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        doc = Document()
        for line in content.split('\n'):
            if line.strip():
                doc.add_paragraph(line)
        
        doc.save(output_path)
        return True
    except Exception as e:
        print(f"Text to Word conversion error: {e}")
        return False

def convert_word_to_text(word_path, output_path):
    """Convert Word to text"""
    try:
        doc = Document(word_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            for para in doc.paragraphs:
                if para.text:
                    f.write(para.text + '\n')
        return True
    except Exception as e:
        print(f"Word to Text conversion error: {e}")
        return False

def convert_image_to_image(input_path, output_path, format):
    """Convert image to another format"""
    try:
        image = Image.open(input_path)
        if format == 'jpg' and image.mode == 'RGBA':
            image = image.convert('RGB')
        image.save(output_path)
        return True
    except Exception as e:
        print(f"Image to Image conversion error: {e}")
        return False

SUPPORTED_CONVERSIONS = {
    'pdf_to_word': {'from': 'pdf', 'to': 'word', 'ext': '.docx', 'mime': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'},
    'pdf_to_ppt': {'from': 'pdf', 'to': 'ppt', 'ext': '.pptx', 'mime': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'},
    'pdf_to_jpg': {'from': 'pdf', 'to': 'jpg', 'ext': '.jpg', 'mime': 'image/jpeg'},
    'pdf_to_png': {'from': 'pdf', 'to': 'png', 'ext': '.png', 'mime': 'image/png'},
    'word_to_pdf': {'from': 'word', 'to': 'pdf', 'ext': '.pdf', 'mime': 'application/pdf'},
    'word_to_txt': {'from': 'word', 'to': 'txt', 'ext': '.txt', 'mime': 'text/plain'},
    'ppt_to_pdf': {'from': 'ppt', 'to': 'pdf', 'ext': '.pdf', 'mime': 'application/pdf'},
    'ppt_to_jpg': {'from': 'ppt', 'to': 'jpg', 'ext': '.jpg', 'mime': 'image/jpeg'},
    'jpg_to_pdf': {'from': 'jpg', 'to': 'pdf', 'ext': '.pdf', 'mime': 'application/pdf'},
    'png_to_pdf': {'from': 'png', 'to': 'pdf', 'ext': '.pdf', 'mime': 'application/pdf'},
    'jpg_to_png': {'from': 'jpg', 'to': 'png', 'ext': '.png', 'mime': 'image/png'},
    'png_to_jpg': {'from': 'png', 'to': 'jpg', 'ext': '.jpg', 'mime': 'image/jpeg'},
    'txt_to_pdf': {'from': 'txt', 'to': 'pdf', 'ext': '.pdf', 'mime': 'application/pdf'},
    'txt_to_word': {'from': 'txt', 'to': 'word', 'ext': '.docx', 'mime': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'},
}

def perform_conversion(input_path, from_type, to_type, output_path):
    if from_type == 'pdf' and to_type == 'word':
        return convert_pdf_to_word(input_path, output_path)
    elif from_type == 'pdf' and to_type == 'ppt':
        return convert_pdf_to_ppt(input_path, output_path)
    elif from_type == 'pdf' and to_type in ['jpg', 'png']:
        return convert_pdf_to_image(input_path, output_path, to_type)
    elif from_type == 'word' and to_type == 'pdf':
        return convert_word_to_pdf(input_path, output_path)
    elif from_type == 'word' and to_type == 'txt':
        return convert_word_to_text(input_path, output_path)
    elif from_type == 'ppt' and to_type == 'pdf':
        return convert_ppt_to_pdf(input_path, output_path)
    elif from_type == 'ppt' and to_type == 'jpg':
        temp_pdf = output_path.replace('.jpg', '_temp.pdf')
        if convert_ppt_to_pdf(input_path, temp_pdf):
            success = convert_pdf_to_image(temp_pdf, output_path, 'jpg')
            if os.path.exists(temp_pdf):
                os.remove(temp_pdf)
            return success
        return False
    elif from_type in ['jpg', 'png', 'gif', 'bmp'] and to_type == 'pdf':
        return convert_image_to_pdf(input_path, output_path)
    elif from_type in ['jpg', 'png'] and to_type in ['jpg', 'png']:
        return convert_image_to_image(input_path, output_path, to_type)
    elif from_type == 'txt' and to_type == 'pdf':
        return convert_text_to_pdf(input_path, output_path)
    elif from_type == 'txt' and to_type == 'word':
        return convert_text_to_word(input_path, output_path)
    return False

# ==================== API ROUTES ====================

@app.route("/api/convert", methods=["POST", "OPTIONS"])
@jwt_required()
def convert_file():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        
        if 'file' not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files['file']
        from_format = request.form.get('from_format', '')
        to_format = request.form.get('to_format', '')
        
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
        
        if not from_format or not to_format:
            return jsonify({"error": "Source and target formats are required"}), 400
        
        from_type = from_format.lower()
        to_type = to_format.lower()
        
        conversion_key = f"{from_type}_to_{to_type}"
        if conversion_key not in SUPPORTED_CONVERSIONS:
            return jsonify({"error": f"Conversion from {from_format} to {to_format} is not supported"}), 400
        
        unique_filename = f"{uuid.uuid4()}_{secure_filename(file.filename)}"
        input_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
        file.save(input_path)
        
        conversion_info = SUPPORTED_CONVERSIONS[conversion_key]
        output_filename = f"{uuid.uuid4()}_{Path(file.filename).stem}{conversion_info['ext']}"
        output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
        
        print(f"📁 Converting: {input_path}")
        print(f"📄 From: {from_type} To: {to_type}")
        
        success = perform_conversion(input_path, from_type, to_type, output_path)
        
        if not success:
            if os.path.exists(input_path):
                os.remove(input_path)
            return jsonify({"error": "Conversion failed. Please try again with a different file."}), 500
        
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            if os.path.exists(input_path):
                os.remove(input_path)
            return jsonify({"error": "Conversion failed - output file is empty"}), 500
        
        file_size = os.path.getsize(output_path)
        conversion_record = {
            "user_id": user_id,
            "original_filename": file.filename,
            "converted_filename": output_filename,
            "from_format": from_format,
            "to_format": to_format,
            "file_path": output_path,
            "size": file_size,
            "created_at": datetime.utcnow()
        }
        converted_files_col.insert_one(conversion_record)
        
        if os.path.exists(input_path):
            os.remove(input_path)
        
        print(f"✅ Conversion successful: {output_filename}")
        
        return send_file(
            output_path,
            as_attachment=True,
            download_name=output_filename,
            mimetype=conversion_info['mime']
        )
        
    except Exception as e:
        print(f"File conversion error: {e}")
        traceback.print_exc()
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

@app.route("/api/conversion-history", methods=["GET", "OPTIONS"])
@jwt_required()
def get_conversion_history():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        conversions = list(converted_files_col.find(
            {"user_id": user_id}
        ).sort("created_at", -1).limit(50))
        
        for conv in conversions:
            conv["_id"] = str(conv["_id"])
        
        response = jsonify(conversions)
        return _corsify_actual_response(response)
        
    except Exception as e:
        print(f"Get conversion history error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/supported-conversions", methods=["GET", "OPTIONS"])
def get_supported_conversions():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        conversions = []
        for key, info in SUPPORTED_CONVERSIONS.items():
            conversions.append({
                "from": info['from'],
                "to": info['to'],
                "description": f"{info['from'].upper()} to {info['to'].upper()}"
            })
        
        response = jsonify(conversions)
        return _corsify_actual_response(response)
        
    except Exception as e:
        print(f"Get supported conversions error: {e}")
        return jsonify({"error": "Internal server error"}), 500

# ==================== AUTHENTICATION ROUTES ====================

@app.route("/api/auth/register", methods=["POST", "OPTIONS"])
def register():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        data = request.json
        username = data.get("username", "").strip().lower()
        password = data.get("password", "")
        display_name = data.get("display_name", username)

        if not username or not password:
            return jsonify({"error": "Username and password required"}), 400
        
        if len(password) < 6:
            return jsonify({"error": "Password must be at least 6 characters"}), 400

        if users_col.find_one({"username": username}):
            return jsonify({"error": "Username already exists"}), 409

        hashed = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())

        user_id = users_col.insert_one({
            "username": username,
            "display_name": display_name,
            "password": hashed,
            "created_at": datetime.utcnow()
        }).inserted_id

        token = create_access_token(identity=str(user_id))

        response = jsonify({
            "token": token,
            "user": {
                "id": str(user_id),
                "username": username,
                "display_name": display_name
            }
        })
        return _corsify_actual_response(response), 201
        
    except Exception as e:
        print(f"Registration error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/auth/login", methods=["POST", "OPTIONS"])
def login():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        data = request.json
        username = data.get("username", "").strip().lower()
        password = data.get("password", "")

        user = users_col.find_one({"username": username})
        
        if not user or not bcrypt.checkpw(password.encode('utf-8'), user["password"]):
            return jsonify({"error": "Invalid username or password"}), 401

        token = create_access_token(identity=str(user["_id"]))

        response = jsonify({
            "token": token,
            "user": {
                "id": str(user["_id"]),
                "username": user["username"],
                "display_name": user.get("display_name", user["username"])
            }
        })
        return _corsify_actual_response(response)
        
    except Exception as e:
        print(f"Login error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/auth/me", methods=["GET", "OPTIONS"])
@jwt_required()
def me():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        user = users_col.find_one({"_id": ObjectId(user_id)})
        
        if not user:
            return jsonify({"error": "User not found"}), 404

        response = jsonify({
            "id": str(user["_id"]),
            "username": user["username"],
            "display_name": user.get("display_name", user["username"])
        })
        return _corsify_actual_response(response)
        
    except Exception as e:
        print(f"Me error: {e}")
        return jsonify({"error": "Internal server error"}), 500

# ==================== CHAT ROUTES ====================

@app.route("/api/chats", methods=["GET", "OPTIONS"])
@jwt_required()
def get_chats():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        chats = list(chats_col.find(
            {"user_id": user_id},
            {"title": 1, "created_at": 1, "updated_at": 1, "model": 1}
        ).sort("updated_at", -1))
        
        response = jsonify([serialize(chat) for chat in chats])
        return _corsify_actual_response(response)
        
    except Exception as e:
        print(f"Get chats error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/chats", methods=["POST", "OPTIONS"])
@jwt_required()
def create_chat():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        data = request.json or {}
        
        now = datetime.utcnow()
        chat_id = chats_col.insert_one({
            "user_id": user_id,
            "title": data.get("title", "New Chat"),
            "model": data.get("model", DEFAULT_MODEL),
            "created_at": now,
            "updated_at": now
        }).inserted_id
        
        response = jsonify(serialize(chats_col.find_one({"_id": chat_id})))
        return _corsify_actual_response(response), 201
        
    except Exception as e:
        print(f"Create chat error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/chats/<chat_id>", methods=["PATCH", "OPTIONS"])
@jwt_required()
def update_chat(chat_id):
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        data = request.json
        
        if "title" in data:
            result = chats_col.update_one(
                {"_id": ObjectId(chat_id), "user_id": user_id},
                {"$set": {"title": data["title"], "updated_at": datetime.utcnow()}}
            )
            
            if result.matched_count == 0:
                return jsonify({"error": "Not found"}), 404
            
            response = jsonify({"message": "Chat updated"})
            return _corsify_actual_response(response)
        
        return jsonify({"error": "No fields to update"}), 400
        
    except Exception as e:
        print(f"Update chat error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/chats/<chat_id>", methods=["DELETE", "OPTIONS"])
@jwt_required()
def delete_chat(chat_id):
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        result = chats_col.delete_one({"_id": ObjectId(chat_id), "user_id": user_id})
        if result.deleted_count == 0:
            return jsonify({"error": "Not found"}), 404
        
        messages_col.delete_many({"chat_id": chat_id})
        
        response = jsonify({"message": "Chat deleted"})
        return _corsify_actual_response(response)
        
    except Exception as e:
        print(f"Delete chat error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/chats/<chat_id>/messages", methods=["GET", "OPTIONS"])
@jwt_required()
def get_messages(chat_id):
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        chat = chats_col.find_one({"_id": ObjectId(chat_id), "user_id": user_id})
        if not chat:
            return jsonify({"error": "Not found"}), 404
        
        messages = list(messages_col.find({"chat_id": chat_id}).sort("created_at", 1))
        response = jsonify([serialize(msg) for msg in messages])
        return _corsify_actual_response(response)
        
    except Exception as e:
        print(f"Get messages error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/chats/<chat_id>/messages/<message_id>", methods=["DELETE", "OPTIONS"])
@jwt_required()
def delete_message(chat_id, message_id):
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        
        chat = chats_col.find_one({"_id": ObjectId(chat_id), "user_id": user_id})
        if not chat:
            return jsonify({"error": "Chat not found"}), 404
        
        message_to_delete = messages_col.find_one({"_id": ObjectId(message_id), "chat_id": chat_id})
        if not message_to_delete:
            return jsonify({"error": "Message not found"}), 404
        
        all_messages = list(messages_col.find({"chat_id": chat_id}).sort("created_at", 1))
        
        message_index = None
        for i, msg in enumerate(all_messages):
            if str(msg["_id"]) == message_id:
                message_index = i
                break
        
        if message_index is None:
            return jsonify({"error": "Message not found in sequence"}), 404
        
        messages_to_delete = []
        
        if message_to_delete["role"] == "user":
            messages_to_delete.append(all_messages[message_index])
            if message_index + 1 < len(all_messages) and all_messages[message_index + 1]["role"] == "assistant":
                messages_to_delete.append(all_messages[message_index + 1])
        else:
            messages_to_delete.append(all_messages[message_index])
        
        file_ids_to_delete = []
        for msg in messages_to_delete:
            if msg.get("file_info") and msg["file_info"].get("file_id"):
                file_ids_to_delete.append(ObjectId(msg["file_info"]["file_id"]))
        
        deleted_ids = [str(msg["_id"]) for msg in messages_to_delete]
        result = messages_col.delete_many({
            "_id": {"$in": [ObjectId(id) for id in deleted_ids]}
        })
        
        if file_ids_to_delete:
            files_col.delete_many({"_id": {"$in": file_ids_to_delete}})
        
        if result.deleted_count == 0:
            return jsonify({"error": "Failed to delete messages"}), 404
        
        response = jsonify({
            "message": f"Deleted {result.deleted_count} message(s)",
            "deleted_count": result.deleted_count,
            "deleted_ids": deleted_ids
        })
        return _corsify_actual_response(response)
        
    except Exception as e:
        print(f"Delete message error: {e}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error"}), 500

# ==================== FILE UPLOAD ROUTES ====================

@app.route("/api/upload", methods=["POST", "OPTIONS"])
@jwt_required()
def upload_file():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        
        if 'file' not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files['file']
        chat_id = request.form.get('chat_id')
        
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
        
        file_ext = file.filename.split('.')[-1] if '.' in file.filename else ''
        unique_filename = f"{uuid.uuid4()}.{file_ext}" if file_ext else str(uuid.uuid4())
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
        
        file.save(file_path)
        file_size = os.path.getsize(file_path)
        
        content_type = file.content_type or ''
        extracted_text = ""
        image_data = None
        
        if content_type.startswith('image/'):
            extracted_text = extract_text_from_image(file_path)
            with open(file_path, 'rb') as f:
                img_base64 = base64.b64encode(f.read()).decode()
                image_data = f"data:{content_type};base64,{img_base64}"
        else:
            extracted_text = extract_full_text_from_file(file_path, file.filename)
        
        file_record = {
            "user_id": user_id,
            "chat_id": chat_id,
            "original_filename": file.filename,
            "stored_filename": unique_filename,
            "file_path": file_path,
            "content_type": content_type,
            "size": file_size,
            "extracted_text": extracted_text,
            "uploaded_at": datetime.utcnow()
        }
        
        if image_data:
            file_record["image_data"] = image_data
        
        file_id = files_col.insert_one(file_record).inserted_id
        
        response_data = {
            "file_id": str(file_id),
            "filename": file.filename,
            "content_type": content_type,
            "size": file_size,
            "extracted_text": extracted_text[:500] if extracted_text else "",
            "file_info": {
                "file_id": str(file_id),
                "filename": file.filename,
                "size": file_size,
                "is_image": content_type.startswith('image/')
            }
        }
        
        if image_data:
            response_data["image_data"] = image_data
            response_data["is_image"] = True
        
        print(f"✅ File uploaded: {file.filename}")
        
        response = jsonify(response_data)
        return _corsify_actual_response(response)
        
    except Exception as e:
        print(f"Upload error: {e}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error"}), 500

def extract_text_from_image(image_path):
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image)
        return text.strip() if text else ""
    except Exception as e:
        print(f"OCR error: {e}")
        return ""

def extract_full_text_from_file(file_path, filename):
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    
    try:
        if ext in ['txt', 'py', 'js', 'html', 'css', 'json', 'md']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext == 'docx':
            doc = Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            return '\n'.join(text)
        elif ext == 'pdf':
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                text = []
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
                return '\n'.join(text)
        else:
            return f"File uploaded: {filename}"
    except Exception as e:
        return f"Error processing file: {str(e)}"

# ==================== DOCUMENT GENERATION ROUTE ====================

@app.route("/api/generate-document", methods=["POST", "OPTIONS"])
@jwt_required()
def generate_document():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        data = request.json
        
        content = data.get("content", "")
        title = data.get("title", "Generated Document")
        author = data.get("author", "LeadSOC-AI")
        
        if not content:
            return jsonify({"error": "No content provided"}), 400
        
        doc = Document()
        
        sections = doc.sections
        for section in sections:
            section.top_margin = Inches(1)
            section.bottom_margin = Inches(1)
            section.left_margin = Inches(1)
            section.right_margin = Inches(1)
        
        title_paragraph = doc.add_heading(title, 0)
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in title_paragraph.runs:
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 122, 204)
        
        author_paragraph = doc.add_paragraph()
        author_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        author_run = author_paragraph.add_run(f"Generated by {author} on {datetime.now().strftime('%B %d, %Y')}")
        author_run.font.size = Pt(12)
        author_run.font.italic = True
        
        doc.add_paragraph()
        
        for line in content.split('\n'):
            if line.strip():
                p = doc.add_paragraph(line)
                p.paragraph_format.line_spacing = 1.5
                p.paragraph_format.space_after = Pt(12)
        
        safe_title = re.sub(r'[^\w\s-]', '', title).strip()
        safe_title = safe_title.replace(' ', '_')
        filename = f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        file_path = os.path.join(app.config['DOCUMENTS_FOLDER'], filename)
        
        doc.save(file_path)
        
        documents_col.insert_one({
            "user_id": user_id,
            "title": title,
            "filename": filename,
            "file_path": file_path,
            "content_preview": content[:500],
            "created_at": datetime.utcnow()
        })
        
        return send_file(
            file_path,
            as_attachment=True,
            download_name=filename,
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        
    except Exception as e:
        print(f"Document generation error: {e}")
        traceback.print_exc()
        return jsonify({"error": f"Failed to generate document: {str(e)}"}), 500

# ==================== STREAMING MESSAGE ROUTE ====================

@app.route("/api/chats/<chat_id>/messages", methods=["POST", "OPTIONS"])
@jwt_required()
def send_message(chat_id):
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        user_id = get_jwt_identity()
        chat = chats_col.find_one({"_id": ObjectId(chat_id), "user_id": user_id})
        if not chat:
            return jsonify({"error": "Not found"}), 404
        
        data = request.json
        user_prompt = data.get("content", "").strip()
        uploaded_file_info = data.get("file_info", {})
        
        if not user_prompt and not uploaded_file_info:
            return jsonify({"error": "Content or file required"}), 400
        
        file_content = ""
        file_record = None
        
        if uploaded_file_info and uploaded_file_info.get("file_id"):
            file_id = uploaded_file_info["file_id"]
            file_record = files_col.find_one({"_id": ObjectId(file_id)})
            if file_record:
                file_content = file_record.get("extracted_text", "")
        
        message_data = {
            "chat_id": chat_id,
            "role": "user",
            "content": user_prompt,
            "created_at": datetime.utcnow()
        }
        
        if uploaded_file_info:
            message_data["file_info"] = uploaded_file_info
        
        messages_col.insert_one(message_data)
        
        def generate():
            full_response = ""
            
            try:
                ollama_messages = []
                
                system_message = "You are LeadSOC-AI, an advanced AI assistant. Help users with questions, document creation, file analysis, and general conversation."
                ollama_messages.append({"role": "system", "content": system_message})
                
                recent_messages = list(messages_col.find({"chat_id": chat_id}).sort("created_at", -1).limit(10))
                recent_messages.reverse()
                
                for msg in recent_messages:
                    if msg["_id"] != message_data["_id"]:
                        role = "user" if msg["role"] == "user" else "assistant"
                        content = msg.get("content", "")
                        if content:
                            ollama_messages.append({"role": role, "content": content})
                
                combined_prompt = user_prompt
                if file_content:
                    if user_prompt:
                        combined_prompt = f"{user_prompt}\n\nFile Content:\n{file_content}"
                    else:
                        combined_prompt = f"Here is the file content:\n\n{file_content}"
                
                if combined_prompt:
                    ollama_messages.append({"role": "user", "content": combined_prompt})
                
                resp = requests.post(
                    f"{OLLAMA_URL}/api/chat",
                    json={"model": DEFAULT_MODEL, "messages": ollama_messages, "stream": True},
                    stream=True,
                    timeout=120
                )
                
                if resp.status_code == 200:
                    for line in resp.iter_lines():
                        if line:
                            try:
                                line_str = line.decode('utf-8')
                                if line_str.startswith('data: '):
                                    line_str = line_str[6:]
                                chunk = json.loads(line_str)
                                if chunk.get("message", {}).get("content"):
                                    token = chunk["message"]["content"]
                                    full_response += token
                                    yield f"data: {json.dumps({'token': token})}\n\n"
                                elif chunk.get("done"):
                                    break
                            except json.JSONDecodeError:
                                continue
                else:
                    if file_content:
                        full_response = f"Here's the file content:\n\n{file_content}"
                    else:
                        full_response = "I'm here to help! How can I assist you?"
                    yield f"data: {json.dumps({'token': full_response})}\n\n"
                        
            except Exception as e:
                print(f"Error: {e}")
                if file_content:
                    full_response = f"Here's the file content:\n\n{file_content}"
                else:
                    full_response = "I'm here to help!"
                yield f"data: {json.dumps({'token': full_response})}\n\n"
            
            if full_response:
                messages_col.insert_one({
                    "chat_id": chat_id,
                    "role": "assistant",
                    "content": full_response,
                    "created_at": datetime.utcnow()
                })
            
            yield f"data: {json.dumps({'done': True})}\n\n"
        
        response = Response(
            stream_with_context(generate()),
            mimetype="text/event-stream",
            headers={
                "X-Accel-Buffering": "no",
                "Cache-Control": "no-cache",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Credentials": "true"
            }
        )
        return response
        
    except Exception as e:
        print(f"Send message error: {e}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error"}), 500

# ==================== TEST ROUTES ====================

@app.route("/api/models", methods=["GET", "OPTIONS"])
@jwt_required()
def get_models():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    try:
        resp = requests.get(f"{OLLAMA_URL}/api/tags", timeout=2)
        if resp.status_code == 200:
            models = [m["name"] for m in resp.json().get("models", [])]
            response = jsonify(models if models else [DEFAULT_MODEL])
        else:
            response = jsonify([DEFAULT_MODEL])
        return _corsify_actual_response(response)
    except Exception as e:
        response = jsonify([DEFAULT_MODEL])
        return _corsify_actual_response(response)

@app.route("/api/test", methods=["GET", "OPTIONS"])
def test():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    response = jsonify({
        "status": "ok",
        "message": "API is working!",
        "timestamp": datetime.utcnow().isoformat()
    })
    return _corsify_actual_response(response)

if __name__ == "__main__":
    print("\n" + "="*70)
    print("🚀 Starting LeadSOC-AI Backend Server")
    print("="*70)
    print(f"📡 Local access: http://localhost:5010")
    print(f"🗄️  MongoDB: {MONGODB_URI}")
    print(f"🤖 Ollama: {OLLAMA_URL}")
    print(f"🤖 Default Model: {DEFAULT_MODEL}")
    print(f"🔄 File Converter: Enabled")
    print(f"   - PDF to Word")
    print(f"   - Word to PDF (LibreOffice recommended)")
    print(f"   - PPT to PDF (LibreOffice recommended)")
    print(f"🌐 CORS: Enabled for ports 3000, 3003, 3005")
    print("="*70)
    print("✅ Server is ready!")
    print("="*70 + "\n")
    
    app.run(debug=True, port=5010, host='0.0.0.0')